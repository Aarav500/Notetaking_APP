"""
Image Flashcards module for AI Note System.
Handles generating image-enhanced flashcards from slides/diagrams.
"""

import os
import logging
import json
import tempfile
import base64
from typing import Dict, Any, List, Optional, Tuple, Union
from pathlib import Path
from datetime import datetime
import re

# Setup logging
logger = logging.getLogger("ai_note_system.processing.image_flashcards")

def extract_images_from_pdf(pdf_path: str, output_dir: Optional[str] = None) -> List[Dict[str, Any]]:
    """
    Extract images from a PDF file.
    
    Args:
        pdf_path (str): Path to the PDF file
        output_dir (str, optional): Directory to save extracted images
        
    Returns:
        List[Dict[str, Any]]: List of extracted images with metadata
    """
    logger.info(f"Extracting images from PDF: {pdf_path}")
    
    try:
        import fitz  # PyMuPDF
        
        # Create temporary directory if output_dir not provided
        if output_dir is None:
            temp_dir = tempfile.mkdtemp()
            output_dir = temp_dir
        else:
            os.makedirs(output_dir, exist_ok=True)
        
        # Open the PDF
        pdf_document = fitz.open(pdf_path)
        
        # Extract images
        images = []
        for page_num, page in enumerate(pdf_document):
            image_list = page.get_images(full=True)
            
            for img_index, img in enumerate(image_list):
                # Get image info
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                # Generate filename
                filename = f"page{page_num+1}_img{img_index+1}.{image_ext}"
                filepath = os.path.join(output_dir, filename)
                
                # Save image
                with open(filepath, "wb") as f:
                    f.write(image_bytes)
                
                # Get image dimensions
                width = base_image.get("width", 0)
                height = base_image.get("height", 0)
                
                # Add to list
                images.append({
                    "page": page_num + 1,
                    "index": img_index + 1,
                    "path": filepath,
                    "filename": filename,
                    "extension": image_ext,
                    "width": width,
                    "height": height
                })
        
        logger.info(f"Extracted {len(images)} images from PDF")
        return images
        
    except ImportError:
        logger.error("PyMuPDF (fitz) not installed. Cannot extract images from PDF.")
        return []
        
    except Exception as e:
        logger.error(f"Error extracting images from PDF: {e}")
        return []

def extract_images_from_slides(slides_path: str, output_dir: Optional[str] = None) -> List[Dict[str, Any]]:
    """
    Extract images from presentation slides (PPTX, etc.).
    
    Args:
        slides_path (str): Path to the slides file
        output_dir (str, optional): Directory to save extracted images
        
    Returns:
        List[Dict[str, Any]]: List of extracted images with metadata
    """
    logger.info(f"Extracting images from slides: {slides_path}")
    
    try:
        import pptx
        from PIL import Image
        
        # Create temporary directory if output_dir not provided
        if output_dir is None:
            temp_dir = tempfile.mkdtemp()
            output_dir = temp_dir
        else:
            os.makedirs(output_dir, exist_ok=True)
        
        # Open the presentation
        presentation = pptx.Presentation(slides_path)
        
        # Extract images
        images = []
        for slide_num, slide in enumerate(presentation.slides):
            for shape_num, shape in enumerate(slide.shapes):
                if hasattr(shape, "image"):
                    # Get image info
                    image_bytes = shape.image.blob
                    
                    # Determine image type
                    if image_bytes.startswith(b'\x89PNG'):
                        image_ext = "png"
                    elif image_bytes.startswith(b'\xff\xd8'):
                        image_ext = "jpg"
                    else:
                        image_ext = "img"
                    
                    # Generate filename
                    filename = f"slide{slide_num+1}_img{shape_num+1}.{image_ext}"
                    filepath = os.path.join(output_dir, filename)
                    
                    # Save image
                    with open(filepath, "wb") as f:
                        f.write(image_bytes)
                    
                    # Get image dimensions
                    try:
                        with Image.open(filepath) as img:
                            width, height = img.size
                    except:
                        width, height = 0, 0
                    
                    # Add to list
                    images.append({
                        "slide": slide_num + 1,
                        "index": shape_num + 1,
                        "path": filepath,
                        "filename": filename,
                        "extension": image_ext,
                        "width": width,
                        "height": height
                    })
        
        logger.info(f"Extracted {len(images)} images from slides")
        return images
        
    except ImportError:
        logger.error("python-pptx not installed. Cannot extract images from slides.")
        return []
        
    except Exception as e:
        logger.error(f"Error extracting images from slides: {e}")
        return []

def extract_images_from_video(video_path: str, output_dir: Optional[str] = None, interval: int = 10) -> List[Dict[str, Any]]:
    """
    Extract keyframes from a video file.
    
    Args:
        video_path (str): Path to the video file
        output_dir (str, optional): Directory to save extracted images
        interval (int): Interval in seconds between keyframes
        
    Returns:
        List[Dict[str, Any]]: List of extracted images with metadata
    """
    logger.info(f"Extracting keyframes from video: {video_path}")
    
    try:
        import cv2
        
        # Create temporary directory if output_dir not provided
        if output_dir is None:
            temp_dir = tempfile.mkdtemp()
            output_dir = temp_dir
        else:
            os.makedirs(output_dir, exist_ok=True)
        
        # Open the video
        video = cv2.VideoCapture(video_path)
        
        # Get video properties
        fps = video.get(cv2.CAP_PROP_FPS)
        frame_count = int(video.get(cv2.CAP_PROP_FRAME_COUNT))
        duration = frame_count / fps
        
        # Calculate frame interval
        frame_interval = int(fps * interval)
        
        # Extract keyframes
        images = []
        frame_num = 0
        img_index = 0
        
        while True:
            # Set position
            video.set(cv2.CAP_PROP_POS_FRAMES, frame_num)
            
            # Read frame
            ret, frame = video.read()
            if not ret:
                break
            
            # Generate filename
            timestamp = frame_num / fps
            minutes = int(timestamp / 60)
            seconds = int(timestamp % 60)
            filename = f"frame_{minutes:02d}m{seconds:02d}s.jpg"
            filepath = os.path.join(output_dir, filename)
            
            # Save frame
            cv2.imwrite(filepath, frame)
            
            # Get image dimensions
            height, width = frame.shape[:2]
            
            # Add to list
            images.append({
                "timestamp": timestamp,
                "timestamp_str": f"{minutes:02d}:{seconds:02d}",
                "frame": frame_num,
                "index": img_index + 1,
                "path": filepath,
                "filename": filename,
                "extension": "jpg",
                "width": width,
                "height": height
            })
            
            # Increment counters
            frame_num += frame_interval
            img_index += 1
            
            # Check if we've reached the end
            if frame_num >= frame_count:
                break
        
        # Release video
        video.release()
        
        logger.info(f"Extracted {len(images)} keyframes from video")
        return images
        
    except ImportError:
        logger.error("OpenCV not installed. Cannot extract keyframes from video.")
        return []
        
    except Exception as e:
        logger.error(f"Error extracting keyframes from video: {e}")
        return []

def generate_image_flashcards(
    content: Dict[str, Any],
    images: List[Dict[str, Any]],
    output_dir: str,
    num_cards: int = 10,
    model: str = "gpt-4-vision-preview"
) -> List[Dict[str, Any]]:
    """
    Generate image-enhanced flashcards from content and images.
    
    Args:
        content (Dict[str, Any]): The content (text, summary, etc.)
        images (List[Dict[str, Any]]): List of extracted images
        output_dir (str): Directory to save flashcards
        num_cards (int): Number of flashcards to generate
        model (str): LLM model to use
        
    Returns:
        List[Dict[str, Any]]: List of generated flashcards
    """
    logger.info(f"Generating {num_cards} image-enhanced flashcards")
    
    try:
        # Ensure output directory exists
        os.makedirs(output_dir, exist_ok=True)
        
        # Filter out small images (likely icons, etc.)
        filtered_images = [img for img in images if img.get("width", 0) > 100 and img.get("height", 0) > 100]
        
        # If no images left after filtering, return empty list
        if not filtered_images:
            logger.warning("No suitable images found for flashcards")
            return []
        
        # Limit number of images to process
        max_images = min(len(filtered_images), 20)  # Process at most 20 images
        selected_images = filtered_images[:max_images]
        
        # Prepare content text
        text = content.get("text", "")
        summary = content.get("summary", "")
        title = content.get("title", "Untitled")
        
        # Use summary if available, otherwise use text (truncated if too long)
        content_text = summary if summary else (text[:2000] + "..." if len(text) > 2000 else text)
        
        # Generate flashcards for each image
        flashcards = []
        
        for img in selected_images:
            # Generate question and answer for this image
            qa_pair = generate_qa_for_image(img["path"], content_text, title, model)
            
            if qa_pair:
                # Add image info to the flashcard
                flashcard = {
                    "question": qa_pair["question"],
                    "answer": qa_pair["answer"],
                    "image_path": img["path"],
                    "image_filename": img["filename"]
                }
                
                # Add timestamp if available (for video frames)
                if "timestamp_str" in img:
                    flashcard["timestamp"] = img["timestamp_str"]
                
                # Add page/slide number if available
                if "page" in img:
                    flashcard["page"] = img["page"]
                elif "slide" in img:
                    flashcard["slide"] = img["slide"]
                
                flashcards.append(flashcard)
                
                # Stop if we've reached the desired number of cards
                if len(flashcards) >= num_cards:
                    break
        
        logger.info(f"Generated {len(flashcards)} image-enhanced flashcards")
        
        # Save flashcards to JSON file
        json_path = os.path.join(output_dir, "image_flashcards.json")
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(flashcards, f, ensure_ascii=False, indent=2)
        
        logger.info(f"Saved flashcards to {json_path}")
        
        return flashcards
        
    except Exception as e:
        logger.error(f"Error generating image-enhanced flashcards: {e}")
        return []

def generate_qa_for_image(
    image_path: str,
    content_text: str,
    title: str,
    model: str = "gpt-4-vision-preview"
) -> Optional[Dict[str, str]]:
    """
    Generate a question-answer pair for an image using a vision-language model.
    
    Args:
        image_path (str): Path to the image
        content_text (str): Related text content
        title (str): Title of the content
        model (str): Vision-language model to use
        
    Returns:
        Optional[Dict[str, str]]: Dictionary with question and answer, or None if failed
    """
    try:
        # For GPT-4 Vision
        if model.startswith("gpt-4-vision"):
            return _generate_qa_with_gpt4_vision(image_path, content_text, title)
        # For other models (fallback to text-only)
        else:
            return _generate_qa_without_vision(image_path, content_text, title)
            
    except Exception as e:
        logger.error(f"Error generating Q&A for image: {e}")
        return None

def _generate_qa_with_gpt4_vision(image_path: str, content_text: str, title: str) -> Optional[Dict[str, str]]:
    """
    Generate a question-answer pair using GPT-4 Vision.
    
    Args:
        image_path (str): Path to the image
        content_text (str): Related text content
        title (str): Title of the content
        
    Returns:
        Optional[Dict[str, str]]: Dictionary with question and answer, or None if failed
    """
    try:
        import openai
        import base64
        
        # Read image and convert to base64
        with open(image_path, "rb") as image_file:
            base64_image = base64.b64encode(image_file.read()).decode('utf-8')
        
        # Create prompt
        prompt = f"""
        You are creating an educational flashcard that includes this image.
        
        The content is about: {title}
        
        Related text content:
        {content_text[:500]}...
        
        Based on the image and the related content, create:
        1. A question that requires understanding the image
        2. A comprehensive answer to the question
        
        The question should be specific to what's shown in the image and relevant to the topic.
        The answer should be educational and informative.
        
        Format your response as JSON with 'question' and 'answer' fields.
        """
        
        # Call OpenAI API
        client = openai.OpenAI()
        response = client.chat.completions.create(
            model="gpt-4-vision-preview",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=1000,
            response_format={"type": "json_object"}
        )
        
        # Parse response
        response_text = response.choices[0].message.content
        qa_pair = json.loads(response_text)
        
        return {
            "question": qa_pair.get("question", ""),
            "answer": qa_pair.get("answer", "")
        }
        
    except ImportError:
        logger.error("OpenAI package not installed. Cannot use GPT-4 Vision.")
        return None
        
    except Exception as e:
        logger.error(f"Error generating Q&A with GPT-4 Vision: {e}")
        return None

def _generate_qa_without_vision(image_path: str, content_text: str, title: str) -> Dict[str, str]:
    """
    Generate a question-answer pair without using vision capabilities.
    This is a fallback method that doesn't actually analyze the image.
    
    Args:
        image_path (str): Path to the image (not used)
        content_text (str): Related text content
        title (str): Title of the content
        
    Returns:
        Dict[str, str]: Dictionary with question and answer
    """
    # Extract filename without extension
    filename = os.path.basename(image_path)
    name_without_ext = os.path.splitext(filename)[0]
    
    # Try to extract meaningful information from filename
    # For example, "slide5_diagram2.png" -> "diagram from slide 5"
    image_desc = name_without_ext
    
    # Replace underscores with spaces
    image_desc = image_desc.replace("_", " ")
    
    # Extract slide/page numbers
    slide_match = re.search(r'slide\s*(\d+)', image_desc, re.IGNORECASE)
    page_match = re.search(r'page\s*(\d+)', image_desc, re.IGNORECASE)
    
    if slide_match:
        slide_num = slide_match.group(1)
        image_desc = f"visual from slide {slide_num}"
    elif page_match:
        page_num = page_match.group(1)
        image_desc = f"figure from page {page_num}"
    
    # Generate generic question and answer
    question = f"Explain the {image_desc} related to {title}."
    answer = f"This {image_desc} illustrates concepts related to {title}. Based on the content: {content_text[:200]}..."
    
    return {
        "question": question,
        "answer": answer
    }

def export_flashcards_to_anki(
    flashcards: List[Dict[str, Any]],
    output_path: str,
    deck_name: str = "Image Flashcards",
    include_metadata: bool = True
) -> Dict[str, Any]:
    """
    Export image flashcards to Anki format.
    
    Args:
        flashcards (List[Dict[str, Any]]): List of flashcards
        output_path (str): Path to save the Anki deck file (.apkg)
        deck_name (str): Name of the Anki deck
        include_metadata (bool): Whether to include metadata
        
    Returns:
        Dict[str, Any]: Result of the export operation
    """
    logger.info(f"Exporting {len(flashcards)} flashcards to Anki")
    
    try:
        # Check if required packages are installed
        try:
            import genanki
            GENANKI_AVAILABLE = True
        except ImportError:
            error_msg = "genanki package not installed. Cannot export to Anki."
            logger.error(error_msg)
            return {"success": False, "error": error_msg}
        
        # Create a unique deck ID
        import random
        deck_id = random.randrange(1 << 30, 1 << 31)
        
        # Create the deck
        deck = genanki.Deck(deck_id, deck_name)
        
        # Create model for image flashcards
        model = create_image_flashcard_model()
        
        # Add flashcards to deck
        media_files = []
        
        for flashcard in flashcards:
            # Get flashcard data
            question = flashcard["question"]
            answer = flashcard["answer"]
            image_path = flashcard["image_path"]
            image_filename = flashcard["image_filename"]
            
            # Create note
            note = create_image_flashcard_note(
                question=question,
                answer=answer,
                image_filename=image_filename,
                model=model,
                metadata=flashcard
            )
            
            # Add note to deck
            deck.add_note(note)
            
            # Add image to media files
            media_files.append(image_path)
        
        # Create package
        package = genanki.Package(deck)
        package.media_files = media_files
        
        # Write package to file
        package.write_to_file(output_path)
        
        logger.info(f"Anki deck with image flashcards saved to {output_path}")
        return {"success": True, "output_path": output_path}
        
    except Exception as e:
        error_msg = f"Error exporting image flashcards to Anki: {str(e)}"
        logger.error(error_msg)
        return {"success": False, "error": error_msg}

def create_image_flashcard_model() -> Any:
    """
    Create an Anki note model for image flashcards.
    
    Returns:
        Any: Anki note model
    """
    import genanki
    import random
    
    model_id = random.randrange(1 << 30, 1 << 31)
    
    model = genanki.Model(
        model_id,
        'AI Note System - Image Flashcard',
        fields=[
            {'name': 'Question'},
            {'name': 'Answer'},
            {'name': 'Image'},
            {'name': 'Metadata'},
        ],
        templates=[
            {
                'name': 'Image Flashcard',
                'qfmt': '''
                <div class="card-question">{{Question}}</div>
                <div class="card-image">{{Image}}</div>
                ''',
                'afmt': '''
                <div class="card-question">{{Question}}</div>
                <div class="card-image">{{Image}}</div>
                <hr>
                <div class="card-answer">{{Answer}}</div>
                {{#Metadata}}
                <div class="card-metadata">
                    <hr>
                    <div class="metadata-label">Additional Information:</div>
                    <div class="metadata-content">{{Metadata}}</div>
                </div>
                {{/Metadata}}
                ''',
            },
        ],
        css='''
        .card {
            font-family: Arial, sans-serif;
            font-size: 16px;
            text-align: left;
            color: #333;
            background-color: #f9f9f9;
            padding: 20px;
        }
        .card-question {
            font-size: 18px;
            color: #2c3e50;
            margin-bottom: 15px;
        }
        .card-image {
            text-align: center;
            margin-bottom: 15px;
        }
        .card-image img {
            max-width: 100%;
            max-height: 300px;
            border: 1px solid #ddd;
            border-radius: 4px;
            padding: 5px;
        }
        .card-answer {
            font-size: 16px;
            color: #333;
            margin-bottom: 15px;
        }
        .card-metadata {
            font-size: 12px;
            color: #777;
            margin-top: 15px;
        }
        .metadata-label {
            font-weight: bold;
            margin-bottom: 5px;
        }
        .metadata-content {
            white-space: pre-wrap;
        }
        '''
    )
    
    return model

def create_image_flashcard_note(
    question: str,
    answer: str,
    image_filename: str,
    model: Any,
    metadata: Optional[Dict[str, Any]] = None
) -> Any:
    """
    Create an Anki note for an image flashcard.
    
    Args:
        question (str): The question
        answer (str): The answer
        image_filename (str): Filename of the image
        model (Any): Anki note model
        metadata (Dict[str, Any], optional): Additional metadata
        
    Returns:
        Any: Anki note
    """
    import genanki
    
    # Create HTML for image
    image_html = f'<img src="{image_filename}" alt="Flashcard Image">'
    
    # Format metadata
    metadata_html = ""
    if metadata:
        metadata_items = []
        
        # Add timestamp if available
        if "timestamp" in metadata:
            metadata_items.append(f"Timestamp: {metadata['timestamp']}")
        
        # Add page/slide number if available
        if "page" in metadata:
            metadata_items.append(f"Page: {metadata['page']}")
        elif "slide" in metadata:
            metadata_items.append(f"Slide: {metadata['slide']}")
        
        metadata_html = "\n".join(metadata_items)
    
    # Create note
    note = genanki.Note(
        model=model,
        fields=[question, answer, image_html, metadata_html]
    )
    
    return note

def process_source_for_flashcards(
    source_path: str,
    source_type: str,
    content: Dict[str, Any],
    output_dir: Optional[str] = None,
    num_cards: int = 10,
    model: str = "gpt-4-vision-preview",
    export_anki: bool = True
) -> Dict[str, Any]:
    """
    Process a source file to generate image-enhanced flashcards.
    
    Args:
        source_path (str): Path to the source file
        source_type (str): Type of source (pdf, slides, video)
        content (Dict[str, Any]): Content data (text, summary, etc.)
        output_dir (str, optional): Directory to save output
        num_cards (int): Number of flashcards to generate
        model (str): LLM model to use
        export_anki (bool): Whether to export to Anki format
        
    Returns:
        Dict[str, Any]: Result of the operation
    """
    logger.info(f"Processing {source_type} source for image flashcards: {source_path}")
    
    try:
        # Create output directory if not provided
        if output_dir is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_dir = os.path.join(tempfile.gettempdir(), f"image_flashcards_{timestamp}")
        
        os.makedirs(output_dir, exist_ok=True)
        
        # Extract images based on source type
        images = []
        
        if source_type.lower() == "pdf":
            images = extract_images_from_pdf(source_path, output_dir)
        elif source_type.lower() in ["slides", "pptx", "powerpoint"]:
            images = extract_images_from_slides(source_path, output_dir)
        elif source_type.lower() in ["video", "mp4", "youtube"]:
            images = extract_images_from_video(source_path, output_dir)
        else:
            return {
                "success": False,
                "error": f"Unsupported source type: {source_type}"
            }
        
        # Check if images were extracted
        if not images:
            return {
                "success": False,
                "error": f"No images extracted from {source_type} source"
            }
        
        # Generate flashcards
        flashcards = generate_image_flashcards(
            content=content,
            images=images,
            output_dir=output_dir,
            num_cards=num_cards,
            model=model
        )
        
        # Check if flashcards were generated
        if not flashcards:
            return {
                "success": False,
                "error": "Failed to generate image flashcards"
            }
        
        result = {
            "success": True,
            "flashcards": flashcards,
            "output_dir": output_dir,
            "json_path": os.path.join(output_dir, "image_flashcards.json")
        }
        
        # Export to Anki if requested
        if export_anki:
            anki_path = os.path.join(output_dir, "image_flashcards.apkg")
            anki_result = export_flashcards_to_anki(
                flashcards=flashcards,
                output_path=anki_path,
                deck_name=content.get("title", "Image Flashcards")
            )
            
            result["anki_export"] = anki_result
            
            if anki_result["success"]:
                result["anki_path"] = anki_result["output_path"]
        
        logger.info(f"Successfully processed source for image flashcards")
        return result
        
    except Exception as e:
        error_msg = f"Error processing source for image flashcards: {str(e)}"
        logger.error(error_msg)
        return {"success": False, "error": error_msg}
def generate_zero_shot_flashcards(
    input_path: str,
    input_type: str = "image",
    output_dir: Optional[str] = None,
    num_cards: int = 5,
    model: str = "gpt-4-vision-preview",
    context: Optional[str] = None
) -> List[Dict[str, Any]]:
    """
    Generate flashcards from unstructured data like screenshots or tweets.
    
    Args:
        input_path (str): Path to the input file (image, screenshot, text file)
        input_type (str): Type of input ('image', 'screenshot', 'social_media', 'text')
        output_dir (str, optional): Directory to save flashcards
        num_cards (int): Number of flashcards to generate
        model (str): LLM model to use
        context (str, optional): Additional context to help with flashcard generation
        
    Returns:
        List[Dict[str, Any]]: List of generated flashcards
    """
    logger.info(f"Generating zero-shot flashcards from {input_type}: {input_path}")
    
    try:
        # Create output directory if not provided
        if output_dir is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_dir = os.path.join(tempfile.gettempdir(), f"zero_shot_flashcards_{timestamp}")
        
        os.makedirs(output_dir, exist_ok=True)
        
        # Process input based on type
        if input_type in ["image", "screenshot"]:
            # For images and screenshots, use vision model
            flashcards = _generate_flashcards_from_image(
                image_path=input_path,
                num_cards=num_cards,
                model=model,
                context=context,
                is_screenshot=(input_type == "screenshot")
            )
        elif input_type == "social_media":
            # For social media posts, extract text and use vision if image is present
            flashcards = _generate_flashcards_from_social_media(
                input_path=input_path,
                num_cards=num_cards,
                model=model,
                context=context
            )
        elif input_type == "text":
            # For text files, extract content and generate flashcards
            flashcards = _generate_flashcards_from_text(
                text_path=input_path,
                num_cards=num_cards,
                model=model,
                context=context
            )
        else:
            logger.error(f"Unsupported input type: {input_type}")
            return []
        
        # Save flashcards to JSON file
        if flashcards:
            json_path = os.path.join(output_dir, "zero_shot_flashcards.json")
            with open(json_path, "w", encoding="utf-8") as f:
                json.dump(flashcards, f, ensure_ascii=False, indent=2)
            
            logger.info(f"Saved {len(flashcards)} zero-shot flashcards to {json_path}")
        
        return flashcards
        
    except Exception as e:
        logger.error(f"Error generating zero-shot flashcards: {e}")
        return []

def _generate_flashcards_from_image(
    image_path: str,
    num_cards: int = 5,
    model: str = "gpt-4-vision-preview",
    context: Optional[str] = None,
    is_screenshot: bool = False
) -> List[Dict[str, Any]]:
    """
    Generate flashcards from an image or screenshot.
    
    Args:
        image_path (str): Path to the image
        num_cards (int): Number of flashcards to generate
        model (str): LLM model to use
        context (str, optional): Additional context to help with flashcard generation
        is_screenshot (bool): Whether the image is a screenshot
        
    Returns:
        List[Dict[str, Any]]: List of generated flashcards
    """
    try:
        import openai
        import base64
        
        # Read image and convert to base64
        with open(image_path, "rb") as image_file:
            base64_image = base64.b64encode(image_file.read()).decode('utf-8')
        
        # Create prompt based on whether it's a screenshot or regular image
        if is_screenshot:
            prompt = f"""
            You are analyzing a screenshot that may contain information from a website, app, document, or social media.
            Your task is to extract key information and create {num_cards} educational flashcards.
            
            For each flashcard, create:
            1. A specific question based on the content in the screenshot
            2. A comprehensive answer to the question
            
            The flashcards should cover the most important concepts or facts visible in the screenshot.
            Focus on creating educational value rather than just describing what's in the image.
            """
        else:
            prompt = f"""
            You are analyzing an image that may contain educational content, diagrams, charts, or other information.
            Your task is to extract key information and create {num_cards} educational flashcards.
            
            For each flashcard, create:
            1. A specific question based on the content in the image
            2. A comprehensive answer to the question
            
            The flashcards should cover the most important concepts or facts visible in the image.
            Focus on creating educational value rather than just describing what's in the image.
            """
        
        # Add context if provided
        if context:
            prompt += f"\n\nAdditional context about this content: {context}"
        
        # Call OpenAI API
        client = openai.OpenAI()
        response = client.chat.completions.create(
            model=model,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=2000,
            response_format={"type": "json_object"}
        )
        
        # Parse response
        response_text = response.choices[0].message.content
        response_json = json.loads(response_text)
        
        # Extract flashcards from response
        flashcards = []
        if "flashcards" in response_json:
            for card in response_json["flashcards"]:
                flashcards.append({
                    "question": card.get("question", ""),
                    "answer": card.get("answer", ""),
                    "source_type": "screenshot" if is_screenshot else "image",
                    "source_path": image_path
                })
        else:
            # Fallback if response doesn't have expected format
            for i in range(1, num_cards + 1):
                question_key = f"question_{i}"
                answer_key = f"answer_{i}"
                if question_key in response_json and answer_key in response_json:
                    flashcards.append({
                        "question": response_json[question_key],
                        "answer": response_json[answer_key],
                        "source_type": "screenshot" if is_screenshot else "image",
                        "source_path": image_path
                    })
        
        return flashcards
        
    except Exception as e:
        logger.error(f"Error generating flashcards from image: {e}")
        return []

def _generate_flashcards_from_social_media(
    input_path: str,
    num_cards: int = 5,
    model: str = "gpt-4",
    context: Optional[str] = None
) -> List[Dict[str, Any]]:
    """
    Generate flashcards from a social media post (text file with optional image path).
    
    Args:
        input_path (str): Path to the text file containing the social media post
        num_cards (int): Number of flashcards to generate
        model (str): LLM model to use
        context (str, optional): Additional context to help with flashcard generation
        
    Returns:
        List[Dict[str, Any]]: List of generated flashcards
    """
    try:
        # Read the social media post content
        with open(input_path, "r", encoding="utf-8") as f:
            content = f.read()
        
        # Check if there's an image path in the content
        image_path = None
        lines = content.split("\n")
        for i, line in enumerate(lines):
            if line.startswith("IMAGE:"):
                image_path = line.replace("IMAGE:", "").strip()
                # Remove the image line from content
                lines.pop(i)
                content = "\n".join(lines)
                break
        
        # If there's an image, use vision model
        if image_path and os.path.exists(image_path):
            # Generate flashcards using both text and image
            import openai
            import base64
            
            # Read image and convert to base64
            with open(image_path, "rb") as image_file:
                base64_image = base64.b64encode(image_file.read()).decode('utf-8')
            
            # Create prompt
            prompt = f"""
            You are analyzing a social media post with both text and image.
            
            The post content is:
            {content}
            
            Your task is to extract key information and create {num_cards} educational flashcards.
            
            For each flashcard, create:
            1. A specific question based on the content in the post
            2. A comprehensive answer to the question
            
            The flashcards should cover the most important concepts or facts in the post.
            Focus on creating educational value rather than just describing the post.
            """
            
            # Add context if provided
            if context:
                prompt += f"\n\nAdditional context about this content: {context}"
            
            # Call OpenAI API
            client = openai.OpenAI()
            response = client.chat.completions.create(
                model="gpt-4-vision-preview",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{base64_image}"
                                }
                            }
                        ]
                    }
                ],
                max_tokens=2000,
                response_format={"type": "json_object"}
            )
            
            # Parse response
            response_text = response.choices[0].message.content
            response_json = json.loads(response_text)
            
        else:
            # Text-only social media post
            import openai
            
            # Create prompt
            prompt = f"""
            You are analyzing a social media post.
            
            The post content is:
            {content}
            
            Your task is to extract key information and create {num_cards} educational flashcards.
            
            For each flashcard, create:
            1. A specific question based on the content in the post
            2. A comprehensive answer to the question
            
            The flashcards should cover the most important concepts or facts in the post.
            Focus on creating educational value rather than just describing the post.
            
            Respond with a JSON object containing an array of flashcards with 'question' and 'answer' fields.
            """
            
            # Add context if provided
            if context:
                prompt += f"\n\nAdditional context about this content: {context}"
            
            # Call OpenAI API
            client = openai.OpenAI()
            response = client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "user", "content": prompt}
                ],
                max_tokens=2000,
                response_format={"type": "json_object"}
            )
            
            # Parse response
            response_text = response.choices[0].message.content
            response_json = json.loads(response_text)
        
        # Extract flashcards from response
        flashcards = []
        if "flashcards" in response_json:
            for card in response_json["flashcards"]:
                flashcards.append({
                    "question": card.get("question", ""),
                    "answer": card.get("answer", ""),
                    "source_type": "social_media",
                    "source_path": input_path,
                    "has_image": image_path is not None
                })
        else:
            # Fallback if response doesn't have expected format
            for i in range(1, num_cards + 1):
                question_key = f"question_{i}"
                answer_key = f"answer_{i}"
                if question_key in response_json and answer_key in response_json:
                    flashcards.append({
                        "question": response_json[question_key],
                        "answer": response_json[answer_key],
                        "source_type": "social_media",
                        "source_path": input_path,
                        "has_image": image_path is not None
                    })
        
        return flashcards
        
    except Exception as e:
        logger.error(f"Error generating flashcards from social media: {e}")
        return []

def _generate_flashcards_from_text(
    text_path: str,
    num_cards: int = 5,
    model: str = "gpt-4",
    context: Optional[str] = None
) -> List[Dict[str, Any]]:
    """
    Generate flashcards from a text file.
    
    Args:
        text_path (str): Path to the text file
        num_cards (int): Number of flashcards to generate
        model (str): LLM model to use
        context (str, optional): Additional context to help with flashcard generation
        
    Returns:
        List[Dict[str, Any]]: List of generated flashcards
    """
    try:
        # Read the text content
        with open(text_path, "r", encoding="utf-8") as f:
            content = f.read()
        
        # Create prompt
        prompt = f"""
        You are analyzing a text document.
        
        The document content is:
        {content}
        
        Your task is to extract key information and create {num_cards} educational flashcards.
        
        For each flashcard, create:
        1. A specific question based on the content in the document
        2. A comprehensive answer to the question
        
        The flashcards should cover the most important concepts or facts in the document.
        Focus on creating educational value.
        
        Respond with a JSON object containing an array of flashcards with 'question' and 'answer' fields.
        """
        
        # Add context if provided
        if context:
            prompt += f"\n\nAdditional context about this content: {context}"
        
        # Call OpenAI API
        import openai
        client = openai.OpenAI()
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "user", "content": prompt}
            ],
            max_tokens=2000,
            response_format={"type": "json_object"}
        )
        
        # Parse response
        response_text = response.choices[0].message.content
        response_json = json.loads(response_text)
        
        # Extract flashcards from response
        flashcards = []
        if "flashcards" in response_json:
            for card in response_json["flashcards"]:
                flashcards.append({
                    "question": card.get("question", ""),
                    "answer": card.get("answer", ""),
                    "source_type": "text",
                    "source_path": text_path
                })
        else:
            # Fallback if response doesn't have expected format
            for i in range(1, num_cards + 1):
                question_key = f"question_{i}"
                answer_key = f"answer_{i}"
                if question_key in response_json and answer_key in response_json:
                    flashcards.append({
                        "question": response_json[question_key],
                        "answer": response_json[answer_key],
                        "source_type": "text",
                        "source_path": text_path
                    })
        
        return flashcards
        
    except Exception as e:
        logger.error(f"Error generating flashcards from text: {e}")
        return []